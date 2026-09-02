import logging
import fitz  # PyMuPDF (PDF)
import pytesseract  # OCR for images
import google.generativeai as genai

from docx import Document as DocxDocument  # DOCX
from pptx import Presentation  # PPTX
from PIL import Image  # Images

from langchain_text_splitters import RecursiveCharacterTextSplitter
from django.conf import settings
from ..models import DocumentChunk

logger = logging.getLogger(__name__)

genai.configure(api_key=settings.GOOGLE_API_KEY)

# Batch size for embedding API calls (Gemini supports up to 100 texts per call)
EMBEDDING_BATCH_SIZE = 100

def extract_text_from_file(file_path, mime_type):
    """Extract text content based on file type."""
    pages = []
    try:
        if mime_type == 'application/pdf':
            with fitz.open(file_path) as doc:
                for idx, page in enumerate(doc):
                    pages.append({"text": page.get_text(), "page_number": idx + 1})
        
        elif 'wordprocessing' in mime_type or mime_type == 'application/msword':
            doc = DocxDocument(file_path)
            text = "\n".join([para.text for para in doc.paragraphs])
            pages.append({"text": text, "page_number": 1})
        
        elif mime_type == 'text/plain':
            with open(file_path, 'r', encoding='utf-8') as f:
                pages.append({"text": f.read(), "page_number": 1})

        elif 'presentation' in mime_type or mime_type == 'application/vnd.ms-powerpoint':
            prs = Presentation(file_path)
            for idx, slide in enumerate(prs.slides):
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_texts.append(shape.text)
                pages.append({"text": "\n\n".join(slide_texts), "page_number": idx + 1})

        elif mime_type.startswith('image/'):
            image = Image.open(file_path)
            text = ""
            # Try Tesseract OCR first
            try:
                text = pytesseract.image_to_string(image)
            except Exception as ocr_err:
                logger.warning(f"Tesseract OCR failed: {ocr_err}")
            
            # Fallback: use Gemini Vision if OCR returned nothing useful
            if not text or not text.strip():
                try:
                    logger.info("Using Gemini Vision to extract text from image...")
                    vision_model = genai.GenerativeModel("gemini-2.5-flash")
                    img_for_gemini = Image.open(file_path)
                    vision_response = vision_model.generate_content(
                        [
                            "Extract ALL text visible in this image. If it's a diagram or chart, "
                            "describe its contents in detail. Return only the extracted text/description.",
                            img_for_gemini,
                        ]
                    )
                    text = vision_response.text if vision_response.text else ""
                except Exception as vision_err:
                    logger.error(f"Gemini Vision fallback also failed: {vision_err}")
                    text = ""
            
            if text.strip():
                pages.append({"text": text, "page_number": 1})
                
    except Exception as e:
        logger.error(f"Error extracting text: {e}")
        return None
        
    return pages


def _batch_embed(texts):
    """
    Embed a list of texts in batches using the Gemini embedding API.
    Returns a flat list of embedding vectors (one per input text).
    """
    all_embeddings = []

    for i in range(0, len(texts), EMBEDDING_BATCH_SIZE):
        batch = texts[i : i + EMBEDDING_BATCH_SIZE]
        result = genai.embed_content(
            model="models/gemini-embedding-001",
            content=batch,
            task_type="retrieval_document",
            output_dimensionality=768,
        )

        # result["embedding"] is a list of embeddings when content is a list
        if isinstance(result, dict):
            embeddings = result.get("embedding", [])
        else:
            embeddings = getattr(result, "embedding", [])

        # Normalize: if single text was passed, embedding is a flat list
        if embeddings and not isinstance(embeddings[0], list):
            embeddings = [embeddings]

        all_embeddings.extend(embeddings)

    return all_embeddings


def process_document_for_rag(document, file_path):
    """Chunk document and save embeddings using batched API calls."""
    try:
        ext_map = {
            'pdf': 'application/pdf',
            'docx': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'txt': 'text/plain',
            'pptx': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'ppt': 'application/vnd.ms-powerpoint',
            'jpg': 'image/jpeg',
            'jpeg': 'image/jpeg',
            'png': 'image/png',
            'gif': 'image/gif',
            'webp': 'image/webp',
        }
        mime_type = ext_map.get(document.extension, 'application/octet-stream')
        
        pages_data = extract_text_from_file(file_path, mime_type)
        if not pages_data:
            logger.warning(f"No text extracted for doc {document.id}")
            return

        # Larger chunks = fewer chunks = fewer API calls + better semantic context
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=2000,
            chunk_overlap=300,
            separators=["\n\n", "\n", ".", " ", ""]
        )

        # Phase 1: Split all pages into chunks (fast, no API calls)
        chunk_data = []  # list of (chunk_text, page_number)
        for page_data in pages_data:
            page_text = page_data["text"]
            page_num = page_data["page_number"]
            if not page_text.strip():
                continue

            chunks = text_splitter.split_text(page_text)
            for chunk_text in chunks:
                if chunk_text.strip():
                    chunk_data.append((chunk_text, page_num))

        if not chunk_data:
            logger.warning(f"No non-empty chunks for doc {document.id}")
            return

        # Phase 2: Batch embed all chunks in one (or few) API call(s)
        chunk_texts = [text for text, _ in chunk_data]
        logger.info(f"Embedding {len(chunk_texts)} chunks for doc {document.id} in batch...")
        embeddings = _batch_embed(chunk_texts)

        if len(embeddings) != len(chunk_data):
            raise ValueError(
                f"Embedding count mismatch: got {len(embeddings)} embeddings "
                f"for {len(chunk_data)} chunks"
            )

        # Phase 3: Build ORM objects and bulk insert
        objs = []
        for idx, ((chunk_text, page_num), embedding) in enumerate(zip(chunk_data, embeddings)):
            objs.append(
                DocumentChunk(
                    document=document,
                    content=chunk_text,
                    embedding=list(embedding),
                    chunk_index=idx,
                    page_number=page_num,
                )
            )

        logger.info(f"Generated {len(objs)} chunks for doc {document.id}")

        if objs:
            DocumentChunk.objects.bulk_create(objs)
        
        document.is_processed = True
        document.save()
        logger.info(f"Successfully processed RAG for doc {document.id}")

    except Exception as e:
        logger.error(f"RAG processing failed: {e}", exc_info=True)